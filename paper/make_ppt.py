"""Build a talk deck (16:9) from the Claude-SD paper + figures, with figure
captions and speaker notes."""
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = pathlib.Path(__file__).parent
FIG = HERE / "figures"

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
BLUE = RGBColor(0x1F, 0x77, 0xB4)
GREY = RGBColor(0x44, 0x44, 0x44)
RED = RGBColor(0xC0, 0x2A, 0x2A)
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _tb(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def set_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def caption(slide, text, top=Inches(6.75)):
    tf = _tb(slide, Inches(0.5), top, Inches(12.3), Inches(0.6))
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(13); p.font.italic = True; p.font.color.rgb = GREY
    p.alignment = PP_ALIGN.CENTER


def bar(slide, color=NAVY, h=Inches(0.12)):
    s = slide.shapes.add_shape(1, 0, 0, SW, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()


def title_slide(title, subtitle, author, notes=None):
    s = prs.slides.add_slide(BLANK)
    box = s.shapes.add_shape(1, 0, Inches(2.3), SW, Inches(2.0))
    box.fill.solid(); box.fill.fore_color.rgb = NAVY; box.line.fill.background()
    tf = _tb(s, Inches(0.8), Inches(2.45), Inches(11.7), Inches(1.7))
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tf2 = _tb(s, Inches(0.8), Inches(4.55), Inches(11.7), Inches(0.8))
    p = tf2.paragraphs[0]; p.text = subtitle
    p.font.size = Pt(18); p.font.italic = True; p.font.color.rgb = BLUE
    tf3 = _tb(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.8))
    p = tf3.paragraphs[0]; p.text = author
    p.font.size = Pt(15); p.font.color.rgb = GREY
    set_notes(s, notes)
    return s


def header(slide, title):
    bar(slide)
    tf = _tb(slide, Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.8))
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = NAVY


def bullets_slide(title, bullets, foot=None, notes=None):
    s = prs.slides.add_slide(BLANK); header(s, title)
    tf = _tb(s, Inches(0.7), Inches(1.25), Inches(12.0), Inches(5.3))
    for i, (txt, lvl, *col) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if lvl == 0 else "   – ") + txt
        p.font.size = Pt(22 - lvl*3); p.level = lvl
        p.font.color.rgb = col[0] if col else (NAVY if lvl == 0 else GREY)
        p.font.bold = bool(lvl == 0 and col and col[0] == RED)
        p.space_after = Pt(7)
    if foot:
        ftf = _tb(s, Inches(0.7), Inches(6.85), Inches(12.0), Inches(0.5))
        fp = ftf.paragraphs[0]; fp.text = foot
        fp.font.size = Pt(12); fp.font.italic = True; fp.font.color.rgb = GREY
    set_notes(s, notes)
    return s


def two_image_slide(title, img_left, img_right, cap=None, notes=None):
    s = prs.slides.add_slide(BLANK); header(s, title)
    half = SW / 2
    top = Inches(1.4); maxh = Inches(4.9)
    for k, img in enumerate((img_left, img_right)):
        pic = s.shapes.add_picture(str(img), int(half*k) + Inches(0.3), top,
                                   width=int(half - Inches(0.6)))
        if pic.height > maxh:
            pic.height = maxh; pic.width = int(pic.height * (pic.width/pic.height))
        pic.left = int(half*k + (half - pic.width)/2)
    if cap:
        caption(s, cap, top=Inches(6.6))
    set_notes(s, notes)
    return s


def image_slide(title, img, bullets=None, img_frac=0.62, cap=None, notes=None):
    s = prs.slides.add_slide(BLANK); header(s, title)
    iw = SW * img_frac
    pic = s.shapes.add_picture(str(img), Inches(0.35), Inches(1.25), width=int(iw))
    if pic.height > Inches(5.4):
        pic.height = Inches(5.4); pic.width = int(pic.height * (pic.width/pic.height))
        pic.left = Inches(0.35)
    if bullets:
        tf = _tb(s, iw + Inches(0.6), Inches(1.4), SW - iw - Inches(0.9), Inches(5.4))
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "• " + b; p.font.size = Pt(16); p.font.color.rgb = GREY
            p.space_after = Pt(8)
    if cap:
        caption(s, cap)
    set_notes(s, notes)
    return s


# ---- slides ---------------------------------------------------------------
title_slide(
    "Building & validating a GPU micromagnetic simulator with an AI coding agent",
    "cross-implementation testing exposes a concurrency defect",
    "Chun-Yeol You  ·  DGIST   |   target: npj Computational Materials",
    notes=("Opening. This talk is a case study: we used an AI coding agent to build a complete GPU "
           "micromagnetic simulator from scratch, validated it against the community standards, and — the "
           "punchline — the validation infrastructure itself caught a real GPU concurrency bug. So it is as "
           "much about a verification methodology as about the code. Target venue: npj Computational Materials."))

bullets_slide("Two questions", [
    ("AI coding agents now write substantial scientific software", 0),
    ("Q1  Can an agent build a *validated, performant* simulator — not a toy?", 0, BLUE),
    ("Q2  What verification makes the result trustworthy?", 0, BLUE),
    ("Micromagnetics: mature solvers (OOMMF, mumax3, mumax+), anchored to µMAG standard problems", 1),
    ("Prior AI work *optimized* existing codes; building+validating from scratch is the harder question", 1),
], notes=("Motivation. Two open questions. First, can an agent produce something real — validated and fast — "
          "not a demo? Second, and more subtle: how do we trust it? Micromagnetics is an ideal testbed because "
          "the field has mature reference solvers and the NIST µMAG standard problems to anchor correctness. "
          "Note prior AI-for-science work mostly optimized existing codes; building and validating one from "
          "scratch is harder, and that's what we do here."))

bullets_slide("Claude-SpinDynamics (Claude-SD)", [
    ("C++20 / CUDA core + Python API, built end-to-end with an AI agent (Claude Code)", 0),
    ("Solves LLG / stochastic-LLG on structured grids", 1),
    ("Distinctive capabilities", 0, BLUE),
    ("float32 AND float64  — unique among GPU micromagnetic codes (mumax family is f32-only)", 1),
    ("two demag FFT backends: cuFFT and VkFFT", 1),
    ("native SOT / STT / Zhang-Li, DMI, RKKY, magnetoelastic; per-cell materials", 1),
    ("RK4 / RK45-DP / Heun + Relax/Minimize; auto-integrator selection", 1),
], notes=("What we built. A full C++20/CUDA simulator with a Python API, solving the Landau-Lifshitz-Gilbert "
          "and stochastic-LLG equations. The headline differentiator versus the mumax family is double "
          "precision — they are float32-only — plus a second FFT backend (VkFFT) and a broad native physics "
          "set: spin-orbit and spin-transfer torques, DMI, RKKY, magnetoelastic, per-cell materials. Three "
          "integrators plus energy minimizers, with an auto-selection helper."))

image_slide("Agent-driven development & architecture", FIG / "Fig1_pipeline.png", [
    "Spec (CLAUDE.md): layers, SI conventions, test map",
    "Loop: kernel + CPU reference + Catch2 test",
    "Accept only if all 4 CUDA builds pass",
    "Fail -> diagnose -> fix (e.g. DMI race)",
], img_frac=0.6,
    cap="Fig 1 — (a) AI-agent development & verification loop; (b) Claude-SD architecture and capabilities.",
    notes=("How it was built. A root specification file fixes the architecture layering, SI conventions, and "
           "the test map. The agent implements each feature as a CUDA kernel plus a CPU reference plus a "
           "Catch2 unit test, and a change is accepted only when all four CUDA build variants pass — f32/f64 "
           "times cuFFT/VkFFT. That redundant build matrix matters later. The loop closes with diagnose-and-"
           "fix; this is exactly how the DMI race we'll discuss was caught and repaired."))

image_slide("Development is quantified", FIG / "Fig6_dev_metrics.png", [
    "155 commits over 31 days",
    "+136k lines of code churn",
    "tests 7 -> 345 (test-driven)",
    "test : source ratio = 0.55",
    "25 effective-field implementations",
], img_frac=0.62,
    cap="Fig 6 — Agent-driven metrics: test-case growth, code composition, commit cadence.",
    notes=("Scale of the effort, mined from git history. 155 commits over a month, ~137k lines of churn. The "
           "key point is the middle panel ratio: a test-to-source ratio of 0.55 and test cases growing in "
           "lockstep with features, from 7 to 345. This is test-driven development, not tests bolted on at "
           "the end — which is what made the cross-checking possible."))

image_slide("Validation: µMAG standard problems SP#1–5", FIG / "Fig2_umag_validation.png", [
    "Agrees with mumax3 / mumax+ / OOMMF within their spread",
    "SP#4 <mx>(1ns): CS −0.979, OOMMF −0.984 (ref −0.986)",
    "SP#1 L_c = 99.7 nm (CS = mumax+)",
    "SP#2 (new): CS = mumax3 to ≤0.006",
], img_frac=0.6,
    cap="Fig 2 — µMAG standard problems SP#1–SP#5 reproduced (Claude-SD, double precision).",
    notes=("Correctness. We reproduce all five µMAG standard problems and agree with the independent solvers "
           "within their mutual spread. For the dynamic switching problem SP#4, our ⟨mx⟩ at 1 ns is −0.979; "
           "the OOMMF double-precision anchor is −0.984; the reference is −0.986 — all within 2%. SP#1 "
           "crossover length matches mumax+, and our newly added SP#2 matches mumax3 to better than 0.006."))

two_image_slide("SP#2 remanence — cross-solver agreement",
                FIG / "FigS3_sp2.png", FIG / "FigS3b_sp2_crosssolver.png",
                cap="Fig S3 — SP#2 remanent ⟨m⟩/M_s and coercivity vs d/ℓ_ex; Claude-SD vs mumax3.",
                notes=("A closer look at SP#2, which we implemented for this work: remanence and coercivity as "
                       "a function of the reduced size d over the exchange length. Claude-SD (left) and the "
                       "cross-solver overlay (right) track mumax3 to within 0.006 across the whole sweep, "
                       "including the dip where flux closure sets in. This is a quantitative, not just "
                       "qualitative, agreement."))

image_slide("Performance: a precision/size trade-off", FIG / "Fig5_landscape.png", [
    "Crossover ~0.1–0.5 M cells",
    "Small / 2-D: CS f32 fastest (5.3× vs mumax3) — CUDA-Graph",
    "Large 3-D: mumax3 / MuMax-CO lead (cuFFT-bound)",
    "f32 = 4–6× f64 (Blackwell Tensor-Core)",
    "No single code wins everywhere",
], img_frac=0.6,
    cap="Fig 5 — Solver capability matrix (left) and throughput landscape (right).",
    notes=("Performance, the honest version. The capability matrix on the left shows where Claude-SD is "
           "broader (double precision, two backends). The landscape on the right shows a crossover near "
           "0.1–0.5 million cells: below it our float32 build with CUDA-Graph replay is fastest — 5.3× over "
           "mumax3 on the SP#4 grid; above it the comparison becomes FFT-bound and the mumax family leads. "
           "Within Claude-SD, float32 is 4–6× faster than float64. The message: no single code wins everywhere."))

two_image_slide("Performance: cross-solver throughput",
                FIG / "Fig3a_throughput.png", FIG / "Fig3b_scenario_bars.png",
                cap="Fig 3 — Throughput vs cell count (crossover) and per-scenario ms per field-eval.",
                notes=("The same story in throughput terms. Left: milliseconds per field-evaluation versus "
                       "cell count for every solver — note we normalize by field-evaluations because the "
                       "codes use different-order integrators, which is the only fair comparison. Right: the "
                       "per-scenario bars. Small/2-D goes to Claude-SD float32; the largest 3-D grids go to "
                       "MuMax-CO, with Claude-SD within about 1.1×."))

bullets_slide("The key result: a hidden defect", [
    ("Relax a seeded skyrmion across the DMI boundary; record topological charge Q", 0),
    ("Same Claude-SD build, repeated identical runs  →  DIFFERENT Q", 0, RED),
    ("run-to-run std of Q up to 0.48  — even in DOUBLE precision", 1, RED),
    ("Q scattered across topological sectors [−1.5, +1.5]", 1),
    ("mumax3 (relax & minimize): deterministic, robust  →  the diagnostic signal", 0, BLUE),
    ("A correct relaxation must not depend on thread scheduling  →  a defect", 1),
], foot="Fig. 4a",
    notes=("Now the core result. We relax a seeded skyrmion across the DMI metastability boundary and record "
           "the topological charge Q. We found that the SAME build, run repeatedly with identical inputs, gave "
           "DIFFERENT Q — scatter up to 0.48 standard deviation, and crucially this happened even in double "
           "precision, so it is not a rounding-tolerance story. Meanwhile mumax3 was perfectly deterministic. "
           "A physically meaningful relaxation must not depend on thread scheduling, so this was a defect in "
           "our code — and the contrast with mumax3 is what flagged it."))

image_slide("Diagnosis → fix → determinism restored", FIG / "Fig4_race_fix.png", [
    "Bisection: only damped-LLG + DMI scattered",
    "Cause: DMI GPU field missing set_stream override",
    "-> ran on its own stream, racing on d_H_out",
    "(FieldSumGPU single-stream mode skips sync)",
    "Fix: add override -> std(Q) = 0, all builds",
    "Invisible to single build & standard problems",
], img_frac=0.6,
    cap="Fig 4 — Cross-validation exposes & fixes a GPU stream race: Q before (scatter) → after (deterministic).",
    notes=("Diagnosis by bisection: every field alone was deterministic, fixed-step RK4 was deterministic, "
           "but damped-LLG relaxation with DMI present scattered. The cause: the GPU DMI field carried its own "
           "CUDA stream but did not override the compositor's set_stream hook, so in single-stream mode — where "
           "the compositor skips inter-field synchronization — DMI ran concurrently and raced on the shared "
           "effective-field buffer. Near the bifurcation, last-bit differences flipped the topological sector. "
           "Adding the override collapses the Q standard deviation to zero across all builds, with all 113 GPU "
           "tests still green. The defect was invisible to any single build and to the standard problems."))

bullets_slide("Why it matters", [
    ("The agent did NOT write defect-free code — it wrote a subtle GPU race", 0),
    ("…that passed every unit test and every standard problem", 1),
    ("What caught it: deliberately redundant design", 0, BLUE),
    ("multiple precisions × FFT backends × independent solvers", 1),
    ("Discrepancy (CS scatter vs mumax3 determinism) was the diagnostic", 1),
    ("Takeaway", 0, RED),
    ("for AI-built scientific software, cross-implementation validation must be a", 1),
    ("first-class design goal — the failure mode is plausible code that is subtly wrong", 1),
], notes=("The generalizable lesson. The agent did not write perfect code — it wrote a subtle concurrency bug "
          "that passed every unit test and every standard problem. What caught it was deliberately redundant "
          "design: multiple precisions, multiple FFT backends, and benchmarking against independent solvers. "
          "Discrepancy was the diagnostic. So the takeaway for AI-assisted scientific software is that cross-"
          "implementation validation should be a first-class design goal, precisely because the failure mode "
          "of agent code is plausible-looking code that is quietly wrong."))

bullets_slide("Summary", [
    ("An AI agent built a µMAG-validated, dual-precision GPU micromagnetic simulator", 0),
    ("Competitive performance; uniquely offers double precision + two FFT backends", 0),
    ("Its multi-build / multi-solver design exposed AND fixed a real GPU concurrency bug", 0, BLUE),
    ("Open & reproducible", 0),
    ("GPLv3 · github.com/mirryou-maker/Claude-SpinDynamics · Zenodo DOI (on release)", 1),
    ("make_report.py regenerates every table & figure from one results file", 1),
], notes=("To summarize: an AI agent built a validated, dual-precision GPU micromagnetic simulator with "
          "competitive performance and a unique capability set. More importantly, its redundant design "
          "exposed and let us fix a real GPU concurrency defect. Everything is open under GPLv3 and fully "
          "reproducible — one command regenerates every table and figure — and it will be archived with a "
          "citable DOI on release. Thank you."))

out = HERE / "Claude-SD_talk.pptx"
prs.save(str(out))
print("wrote", out, "(", len(prs.slides._sldIdLst), "slides )")
